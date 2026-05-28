"""
Renderer de PowerPoint para fichas hospitalarias.
- Reemplazo por NOMBRE de shape (no por contenido). Cada shape de la plantilla
  tiene un nombre único (TXT_*, TABLA_*, FOTO_*) y mapeamos directo desde ahí.
- Si un dato viene None / vacío / 0 cuando no debería, dejamos el placeholder
  original de la plantilla intacto (XXX, AA, etc.) para que sea evidente que falta.
- Preservamos el formato del primer run de cada párrafo/celda al escribir.
"""

from pathlib import Path
from typing import Optional, Dict, Any, List
from copy import deepcopy

from pptx import Presentation
from pptx.util import Emu
from ficha.src.core.data_models import FichaData
from ficha.src.core.config import AppConfig


# HELPERS DE FORMATO
# Estrategia de preservación de formato:
# Un párrafo en un .pptx puede tener su formato definido en dos lugares:
#   1. dentro de un <a:r><a:rPr.../></a:r>  (cuando ya hay texto)
#   2. en <a:endParaRPr/> al final del párrafo (cuando el párrafo está vacío;
#      es el formato que tendría el siguiente texto que se escriba ahí)
#
# Si solo capturamos el caso 1, perdemos el formato de placeholders vacíos
# como TXT_ESPECIALIDADES o las celdas de número de las tablas.
#
# La estrategia de "clonar el rPr/endParaRPr crudo del XML" es más robusta
# que extraer propiedad por propiedad: hereda todo (color, fuente, tamaño,
# negrita, espaciado de letras, kerning, etc.) sin tener que enumerar cada
# propiedad explícitamente.

NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
TAG_RPR = f"{{{NS_A}}}rPr"
TAG_END_PARA_RPR = f"{{{NS_A}}}endParaRPr"
TAG_PPR = f"{{{NS_A}}}pPr"
TAG_R = f"{{{NS_A}}}r"
TAG_T = f"{{{NS_A}}}t"


def _extraer_rPr_template(p_element) -> Optional[object]:
    """
    Devuelve un nodo XML <a:rPr> que representa el formato a aplicar al
    siguiente texto que se inserte en este párrafo.

    Orden de búsqueda:
      1. rPr del primer run con texto del párrafo.
      2. endParaRPr del párrafo (renombrado a rPr en el clon).

    Devuelve None si no encuentra nada.
    """
    # 1. Buscar primer run con rPr explícito
    for r_node in p_element.findall(TAG_R):
        rPr_node = r_node.find(TAG_RPR)
        if rPr_node is not None:
            return deepcopy(rPr_node)

    # 2. Fallback a endParaRPr — lo clonamos y lo renombramos a rPr
    end_rpr = p_element.find(TAG_END_PARA_RPR)
    if end_rpr is not None:
        clon = deepcopy(end_rpr)
        clon.tag = TAG_RPR
        return clon

    return None


def _crear_run_con_formato(p_element, texto: str, rPr_template):
    """
    Crea un nuevo <a:r> dentro del párrafo con el formato dado y el texto.
    Si rPr_template es None, crea el run sin formato explícito (heredará
    del master). Devuelve el nodo <a:r> creado.
    """
    from lxml import etree

    r_new = etree.SubElement(p_element, TAG_R)
    if rPr_template is not None:
        r_new.append(deepcopy(rPr_template))
    t_new = etree.SubElement(r_new, TAG_T)
    t_new.text = texto
    return r_new


def _reescribir_parrafo(p_element, texto: str):
    """
    Reescribe el contenido textual de un párrafo preservando su formato:
      1. Captura el rPr/endParaRPr existente como template.
      2. Borra todos los runs y endParaRPr.
      3. Crea un único run nuevo con el texto y el rPr clonado.
      4. Conserva el pPr (alineación, sangría, bullets) intacto.
    """
    rPr_template = _extraer_rPr_template(p_element)

    # Borrar runs y endParaRPr existentes; conservar pPr
    for child in list(p_element):
        tag = child.tag
        if tag == TAG_PPR:
            continue
        p_element.remove(child)

    _crear_run_con_formato(p_element, texto, rPr_template)

# RENDERER PRINCIPAL
class PPTRenderer:
    """Genera el PPTX final a partir de FichaData + dict de fotos."""

    def __init__(self, config: AppConfig):
        self.config = config
        self.ruta_salida = config.ruta_salida

    def generar_ppt(self, ficha: FichaData, fotos: dict) -> Path:
        ruta_plantilla = (self.config.plantilla_hr_pptx
                          if ficha.es_hr
                          else self.config.plantilla_pptx)
        if not ruta_plantilla.exists():
            raise FileNotFoundError(f"No se encontró la plantilla: {ruta_plantilla}")

        prs = Presentation(ruta_plantilla)

        textos = self._construir_mapa_textos(ficha)
        listas = self._construir_mapa_listas(ficha)
        tablas = self._construir_mapa_tablas(ficha)
        fotos_map = {
            "FOTO_FACHADA": fotos.get("fachada"),
            "FOTO_1":       fotos.get("interior_1"),
            "FOTO_2":       fotos.get("interior_2"),
            "FOTO_3":       fotos.get("interior_3"),
            "FOTO_4":       fotos.get("interior_4"),
        }

        for slide in prs.slides:
            # Iteramos sobre una copia de la lista de shapes porque podemos
            # remover algunos (al insertar fotos).
            for shape in list(slide.shapes):
                self._procesar_shape(slide, shape, textos, listas, tablas, fotos_map)

        ruta_final = self._ruta_salida_para(ficha)
        prs.save(ruta_final)
        return ruta_final

    # shapes
    def _construir_mapa_textos(self, ficha: FichaData) -> Dict[str, Optional[str]]:
        """
        Mapa nombre_shape → valor a escribir.
        Si el valor es None o cadena vacía, se DEJA el placeholder de la
        plantilla intacto (no escribimos nada).
        """
        def opcional(valor, formato="{}"):
            """Devuelve None si el valor es 0 / None / NaN para dejar placeholder."""
            if valor is None:
                return None
            try:
                if isinstance(valor, (int, float)) and valor == 0:
                    return None
            except Exception:
                pass
            return formato.format(valor)

        anio_inicio = ficha.anio_inicio if ficha.anio_inicio != "N/D" else None

        return {
            # Diapo 1 y 2 — encabezados
            "TXT_NOMBRE_UNIDAD":    ficha.nombre.upper() if ficha.nombre else None,
            "TXT_ESTADO":           ficha.estado.upper() if ficha.estado else None,

            # Diapo 1 — generales
            "TXT_ANTIGUEDAD":       opcional(ficha.antiguedad, "{:d}"),
            "TXT_ANIO_INICIO":      f"({anio_inicio})" if anio_inicio else None,
            "TXT_POBLACION":        ficha.poblacion_regionalizacion or None,
            "TXT_CAMAS":            opcional(ficha.total_camas, "{:,}"),
            "TXT_ESPECIALIDADES":   (f"{ficha.total_especialidades:,} especialidades"
                                     if ficha.total_especialidades else None),

            # Diapo 2 — personal
            "TXT_PERSONAL_TOTAL":   opcional(ficha.total_personal, "{:,}"),
            "TXT_TOTAL_MED":        opcional(ficha.total_medicos, "{:,}"),
            "TXT_TOTAL_ENF":        opcional(ficha.total_enfermeras, "{:,}"),
            "TXT_TOTAL_OTROS":      opcional(ficha.total_otros, "{:,}"),

            # Diapo 2 — productividad y día típico
            "TXT_DIA_TIPICO_TOTAL": opcional(ficha.dt_total, "{:,}"),
            "TXT_PROD_TOTAL":       opcional(ficha.prod_total, "{:,}"),
            "TXT_DATE_PROD":        ficha.fecha_corte_texto or None,
        }

    def _construir_mapa_listas(self, ficha: FichaData) -> Dict[str, list]:
        return {
            "TXT_LISTA_CAP_INSTALADA": ficha.cartera_servicios_lista,
            "TXT_LISTA_EQUIPO":        ficha.equipo_relevante_lista,
            "TXT_LISTA_PERSONAL":      ficha.top_especialidades_lista,
        }

    def _construir_mapa_tablas(self, ficha: FichaData) -> Dict[str, callable]:
        return {
            "TABLA_DIA_TIPICO":         self._llenar_tabla_dia_tipico,
            "TABLA_PRODUCTIVIDAD":      self._llenar_tabla_productividad,
            "TABLA_PRODUCTIVIDAD_AÑOS": self._llenar_tabla_anios,
            "TABLA_PROD_HR":            self._llenar_tabla_hr,
        }

    # procesamiento por shape
    def _procesar_shape(self, slide, shape, textos, listas, tablas, fotos_map):
        nombre = shape.name

        # 1. ¿Es un placeholder de foto?
        if nombre in fotos_map:
            ruta_foto = fotos_map[nombre]
            if ruta_foto and Path(ruta_foto).exists():
                self._reemplazar_imagen(slide, shape, ruta_foto)
            # Si no hay foto, dejamos el shape de la plantilla intacto.
            return

        # 2. ¿Es una tabla con función específica?
        if shape.has_table and nombre in tablas:
            tablas[nombre](shape.table, fotos_map.get("__ficha__"), shape, slide)  # placeholder, ver abajo
            return

        # 3. ¿Es un shape de texto que es lista?
        if shape.has_text_frame and nombre in listas:
            items = listas[nombre]
            if items:
                self._escribir_lista(shape, items)
            # Si la lista está vacía, dejamos el "Lista" de la plantilla.
            return

        # 4. ¿Es un placeholder de texto simple?
        if shape.has_text_frame and nombre in textos:
            valor = textos[nombre]
            if valor is not None and valor != "":
                self._reemplazar_texto_shape(shape, str(valor))
            return

        # 5. Para cualquier otro shape, no hacemos nada (decoración, headers, etc.)

    # escritura de texto en un shape (preserva formato)
    def _reemplazar_texto_shape(self, shape, nuevo_texto: str):
        """
        Reemplaza el texto del shape por `nuevo_texto`, preservando el
        formato del primer párrafo de la plantilla (incluso si está vacío,
        gracias a que leemos endParaRPr).
        """
        tf = shape.text_frame
        if not tf.paragraphs:
            return

        primer_p = tf.paragraphs[0]._p

        # Borrar todos los párrafos posteriores del XML
        for p in list(tf.paragraphs[1:]):
            p._p.getparent().remove(p._p)

        _reescribir_parrafo(primer_p, nuevo_texto)

    def _escribir_lista(self, shape, items: List[str]):
        """
        Escribe una lista de strings, un párrafo por ítem.
        Preserva tanto el formato del párrafo (bullets, sangría) como
        el formato del texto (fuente, tamaño, color).
        """
        tf = shape.text_frame
        if not tf.paragraphs or not items:
            return

        primer_p_elem = tf.paragraphs[0]._p

        # Capturar plantilla de pPr y rPr antes de tocar nada
        plantilla_pPr = primer_p_elem.find(TAG_PPR)
        if plantilla_pPr is not None:
            plantilla_pPr = deepcopy(plantilla_pPr)
        plantilla_rPr = _extraer_rPr_template(primer_p_elem)

        # Borrar párrafos posteriores
        for p in list(tf.paragraphs[1:]):
            p._p.getparent().remove(p._p)

        # Reescribir el primer párrafo con el primer ítem
        _reescribir_parrafo(primer_p_elem, str(items[0]))

        # Crear párrafos nuevos para cada ítem adicional, replicando pPr/rPr
        for item in items[1:]:
            nuevo_p = tf.add_paragraph()
            nuevo_p_elem = nuevo_p._p

            # Borrar el pPr default que python-pptx pudo haber creado
            for child in list(nuevo_p_elem):
                nuevo_p_elem.remove(child)

            # Insertar el pPr de la plantilla (si lo había)
            if plantilla_pPr is not None:
                nuevo_p_elem.append(deepcopy(plantilla_pPr))

            _crear_run_con_formato(nuevo_p_elem, str(item), plantilla_rPr)

    # inserción de imagen
    def _reemplazar_imagen(self, slide, shape_placeholder, ruta_imagen: Path):
        x, y   = shape_placeholder.left, shape_placeholder.top
        cx, cy = shape_placeholder.width, shape_placeholder.height
        slide.shapes.add_picture(str(ruta_imagen), x, y, cx, cy)
        sp = shape_placeholder._element
        sp.getparent().remove(sp)

    # escritura en celdas de tabla (preserva formato) 
    def _set_cell(self, table, row: int, col: int, value: str):
        """
        Escribe en la celda preservando formato (incluyendo el caso de
        celdas vacías, donde el formato vive en endParaRPr).
        """
        try:
            cell = table.cell(row, col)
        except (IndexError, KeyError):
            return

        tf = cell.text_frame
        if not tf.paragraphs:
            return

        primer_p = tf.paragraphs[0]._p

        # Borrar párrafos extras
        for p in list(tf.paragraphs[1:]):
            p._p.getparent().remove(p._p)

        _reescribir_parrafo(primer_p, str(value))

    # llenadores de tablas (firma uniforme) 
    # Las tablas reciben (table, ficha, shape, slide) por compatibilidad pero
    # solo usan ficha. El procesador les pasa la ficha guardada en fotos_map.
    def _llenar_tabla_dia_tipico(self, table, ficha, shape=None, slide=None):
        """
        TABLA_DIA_TIPICO — 3 filas x 4 cols.
        Layout (las labels de col 1 y 3 ya vienen en la plantilla):
          (0,0)=esp        (0,1)=Consultas Esp     (0,2)=urg      (0,3)=Urgencias
          (1,0)=iqx        (1,1)=Cirugías          (1,2)=mf|partos (1,3)=label MF si aplica
          (2,0)=egresos    (2,1)=Egresos           (2,2)=         (2,3)=
        """
        ficha = self._ficha_actual
        if not ficha:
            return

        if ficha.dt_esp:     self._set_cell(table, 0, 0, f"{ficha.dt_esp:,}")
        if ficha.dt_urg:     self._set_cell(table, 0, 2, f"{ficha.dt_urg:,}")
        if ficha.dt_iqx:     self._set_cell(table, 1, 0, f"{ficha.dt_iqx:,}")
        if ficha.dt_egresos: self._set_cell(table, 2, 0, f"{ficha.dt_egresos:,}")

        if ficha.tiene_mf and ficha.dt_mf:
            self._set_cell(table, 1, 2, f"{ficha.dt_mf:,}")
            self._set_cell(table, 1, 3, "Consultas Medicina Familiar")
        elif ficha.dt_partos:
            self._set_cell(table, 1, 2, f"{ficha.dt_partos:,}")
            self._set_cell(table, 1, 3, "Partos Atendidos")

    def _llenar_tabla_productividad(self, table, ficha, shape=None, slide=None):
        """
        TABLA_PRODUCTIVIDAD (acumulada) — 3 filas x 4 cols.
        (0,0)=acum_esp   (0,1)=Consultas Esp   (0,2)=prod_urg   (0,3)=Urgencias
        (1,0)=acum_iqx   (1,1)=Cirugías        (1,2)=prod_egr   (1,3)=Egresos
        (2,0)=           (2,1)=                (2,2)=acum_mf    (2,3)=label MF
        """
        ficha = self._ficha_actual
        if not ficha:
            return

        if ficha.acum_esp: self._set_cell(table, 0, 0, f"{ficha.acum_esp:,}")
        if ficha.prod_urg: self._set_cell(table, 0, 2, f"{ficha.prod_urg:,}")
        if ficha.acum_iqx: self._set_cell(table, 1, 0, f"{ficha.acum_iqx:,}")
        if ficha.prod_egr: self._set_cell(table, 1, 2, f"{ficha.prod_egr:,}")
        if ficha.tiene_mf and ficha.acum_mf:
            self._set_cell(table, 2, 2, f"{ficha.acum_mf:,}")
            self._set_cell(table, 2, 3, "Consultas Medicina Familiar")

    def _llenar_tabla_anios(self, table, ficha, shape=None, slide=None):
        """
        TABLA_PRODUCTIVIDAD_AÑOS — 3 filas x 6 cols.
        Fila 0: headers (no se toca).
        Col 1: año anterior | Col 2: aumento | Col 3: año actual
        Col 4: meta         | Col 5: % avance
        """
        ficha = self._ficha_actual
        if not ficha:
            return

        # Especialidades (fila 1)
        if ficha.comp_esp_ant:     self._set_cell(table, 1, 1, f"{ficha.comp_esp_ant:,}")
        if ficha.comp_esp_aumento: self._set_cell(table, 1, 2, f"{ficha.comp_esp_aumento:,}")
        if ficha.acum_esp:         self._set_cell(table, 1, 3, f"{ficha.acum_esp:,}")
        if ficha.meta_esp:         self._set_cell(table, 1, 4, f"{ficha.meta_esp:,}")
        if ficha.avance_esp_str and ficha.avance_esp_str != "S/M":
            self._set_cell(table, 1, 5, ficha.avance_esp_str)

        # Cirugías (fila 2)
        if ficha.comp_iqx_ant:     self._set_cell(table, 2, 1, f"{ficha.comp_iqx_ant:,}")
        if ficha.comp_iqx_aumento: self._set_cell(table, 2, 2, f"{ficha.comp_iqx_aumento:,}")
        if ficha.acum_iqx:         self._set_cell(table, 2, 3, f"{ficha.acum_iqx:,}")
        if ficha.meta_iqx:         self._set_cell(table, 2, 4, f"{ficha.meta_iqx:,}")
        if ficha.avance_iqx_str and ficha.avance_iqx_str != "S/M":
            self._set_cell(table, 2, 5, ficha.avance_iqx_str)

    def _llenar_tabla_hr(self, table, ficha, shape=None, slide=None):
        """
        TABLA_PROD_HR — 3 filas x 4 cols (solo HR).
        Fila 0: headers. Col 1: año ant | Col 2: año act | Col 3: variación %
        """
        ficha = self._ficha_actual
        if not ficha:
            return

        if ficha.hr_esp_ant: self._set_cell(table, 1, 1, f"{ficha.hr_esp_ant:,}")
        if ficha.hr_esp_act: self._set_cell(table, 1, 2, f"{ficha.hr_esp_act:,}")
        if ficha.hr_esp_var and ficha.hr_esp_var != "N/D":
            self._set_cell(table, 1, 3, ficha.hr_esp_var)

        if ficha.hr_iqx_ant: self._set_cell(table, 2, 1, f"{ficha.hr_iqx_ant:,}")
        if ficha.hr_iqx_act: self._set_cell(table, 2, 2, f"{ficha.hr_iqx_act:,}")
        if ficha.hr_iqx_var and ficha.hr_iqx_var != "N/D":
            self._set_cell(table, 2, 3, ficha.hr_iqx_var)

    # helpers 

    def _ruta_salida_para(self, ficha: FichaData) -> Path:
        tipo = (ficha.tipo_unidad or "HOSP").strip().replace(" ", "_").replace("/", "-")
        nombre_seguro = (ficha.nombre or "")[:30].strip().replace("/", "-").replace("\\", "-")
        nombre_archivo = f"Ficha_{tipo}_{ficha.clave}_{nombre_seguro}.pptx"
        return self.ruta_salida / nombre_archivo

_orig_generar = PPTRenderer.generar_ppt

def _generar_ppt_wrapped(self, ficha: FichaData, fotos: dict) -> Path:
    self._ficha_actual = ficha
    try:
        return _orig_generar(self, ficha, fotos)
    finally:
        self._ficha_actual = None

PPTRenderer.generar_ppt = _generar_ppt_wrapped