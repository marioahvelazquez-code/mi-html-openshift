### Capa de presentacion de la Ficha Presidencial
### Escribe los valores de una FichaData en la plantilla, buscando los shapes
### por nombre. No contiene reglas de negocio: solo sabe que dato va en que
### shape. Un valor None se omite y conserva el placeholder de la plantilla.

import copy

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

import config


def recolectar_shapes(shapes, mapa=None):
    """Recorre los shapes de la diapositiva, incluidos los de grupos, y los
    devuelve indexados por nombre.
    Parametros:
    shapes: coleccion de shapes a recorrer
    mapa: diccionario acumulador, se crea en la primera llamada"""
    if mapa is None:
        mapa = {}
    for shape in shapes:
        mapa[shape.name] = shape
        if shape.shape_type == 6:
            recolectar_shapes(shape.shapes, mapa)
    return mapa


def eliminar_shape(shape):
    """Quita un shape de la diapositiva.
    Parametros:
    shape: shape de python-pptx"""
    elemento = shape._element
    elemento.getparent().remove(elemento)


def escribir_en_parrafo(parrafo_xml, texto):
    """Escribe el texto en el primer run del parrafo y elimina los demas, de
    modo que se conserve el formato original del run.
    Parametros:
    parrafo_xml: elemento a:p
    texto: contenido a escribir"""
    runs = parrafo_xml.findall(qn("a:r"))
    if not runs:
        return False
    for extra in runs[1:]:
        parrafo_xml.remove(extra)
    nodo_texto = runs[0].find(qn("a:t"))
    if nodo_texto is None:
        return False
    nodo_texto.text = texto
    return True


def clonar_run(parrafo_xml, texto, negritas=False, color=None, proporcion=None):
    """Agrega un run clonado del primero del parrafo, con el texto indicado.
    Clonarlo conserva fuente, tamano y demas atributos del run original.
    Parametros:
    parrafo_xml: elemento a:p que ya tiene al menos un run
    texto: contenido del run nuevo
    negritas: aplica negritas al run nuevo
    color: color en hexadecimal, opcional
    proporcion: factor que se aplica al tamano de fuente, opcional"""
    runs = parrafo_xml.findall(qn("a:r"))
    if not runs:
        return

    nuevo = copy.deepcopy(runs[0])
    nuevo.find(qn("a:t")).text = texto

    propiedades = nuevo.find(qn("a:rPr"))
    if propiedades is None:
        propiedades = nuevo.makeelement(qn("a:rPr"), {})
        nuevo.insert(0, propiedades)

    if negritas:
        propiedades.set("b", "1")

    if proporcion is not None and propiedades.get("sz"):
        propiedades.set("sz", str(int(int(propiedades.get("sz")) * proporcion)))

    if color is not None:
        for hijo in propiedades.findall(qn("a:solidFill")):
            propiedades.remove(hijo)
        relleno = propiedades.makeelement(qn("a:solidFill"), {})
        srgb = propiedades.makeelement(qn("a:srgbClr"), {"val": color})
        relleno.append(srgb)
        propiedades.append(relleno)

    # El run nuevo va justo despues del ultimo run existente. Agregarlo al
    # final del parrafo lo dejaria despues del nodo endParaRPr, que las
    # plantillas exportadas de Google Slides si traen, y PowerPoint descarta
    # cualquier run que quede despues de ese nodo.
    runs[-1].addnext(nuevo)


def dar_formato_run(run_xml, negritas=False, color=None):
    """Aplica negritas y color al run indicado, conservando lo demas.
    Parametros:
    run_xml: elemento a:r
    negritas: aplica negritas
    color: color en hexadecimal, opcional"""
    propiedades = run_xml.find(qn("a:rPr"))
    if propiedades is None:
        propiedades = run_xml.makeelement(qn("a:rPr"), {})
        run_xml.insert(0, propiedades)

    if negritas:
        propiedades.set("b", "1")

    if color is not None:
        for hijo in propiedades.findall(qn("a:solidFill")):
            propiedades.remove(hijo)
        relleno = propiedades.makeelement(qn("a:solidFill"), {})
        relleno.append(propiedades.makeelement(qn("a:srgbClr"), {"val": color}))
        propiedades.append(relleno)


def aplicar_sangria(parrafo_xml, centimetros):
    """Aplica sangria izquierda al parrafo.
    Parametros:
    parrafo_xml: elemento a:p
    centimetros: sangria en centimetros"""
    propiedades = parrafo_xml.find(qn("a:pPr"))
    if propiedades is None:
        propiedades = parrafo_xml.makeelement(qn("a:pPr"), {})
        parrafo_xml.insert(0, propiedades)
    propiedades.set("marL", str(int(Cm(centimetros))))
    propiedades.set("indent", "0")


def escribir_con_sufijo_millones(parrafo_xml, texto):
    """Escribe una cifra abreviada y su sufijo en dos runs, con el sufijo en
    menor tamano.
    Parametros:
    parrafo_xml: elemento a:p
    texto: contenido con el separador de sufijo"""
    cifra, sufijo = texto.split(config.SEPARADOR_SUFIJO, 1)
    if escribir_en_parrafo(parrafo_xml, cifra):
        clonar_run(parrafo_xml, f" {sufijo}", proporcion=config.PROPORCION_SUFIJO)


def escribir_valor(parrafo_xml, valor):
    """Escribe un valor en un parrafo, partiendolo en dos runs cuando trae el
    separador de sufijo de millones.
    Parametros:
    parrafo_xml: elemento a:p
    valor: texto a escribir"""
    if config.SEPARADOR_SUFIJO in valor:
        escribir_con_sufijo_millones(parrafo_xml, valor)
    else:
        escribir_en_parrafo(parrafo_xml, valor)


def escribir_texto(shape, valor):
    """Escribe un valor en un shape de texto conservando su formato.
    Parametros:
    shape: shape con text_frame
    valor: texto a escribir, o None para dejar el placeholder"""
    if valor is None or shape is None or not shape.has_text_frame:
        return
    parrafos = shape.text_frame.paragraphs
    if not parrafos:
        return
    escribir_valor(parrafos[0]._p, valor)
    cuerpo = parrafos[0]._p.getparent()
    for parrafo in parrafos[1:]:
        cuerpo.remove(parrafo._p)


def escribir_listado(shape, renglones):
    """Escribe los renglones de un listado en un shape de texto, clonando el
    formato del parrafo original. Los encabezados de estatus van en negritas y
    con su color; los proyectos llevan sangria.
    Parametros:
    shape: shape con text_frame
    renglones: lista de diccionarios con texto y marca de encabezado"""
    if shape is None or not shape.has_text_frame:
        return
    parrafos = shape.text_frame.paragraphs
    if not parrafos:
        return

    modelo = copy.deepcopy(parrafos[0]._p)
    cuerpo = parrafos[0]._p.getparent()
    for parrafo in parrafos:
        cuerpo.remove(parrafo._p)

    for renglon in renglones:
        nuevo = copy.deepcopy(modelo)
        if not escribir_en_parrafo(nuevo, renglon["texto"]):
            cuerpo.append(nuevo)
            continue

        if renglon.get("encabezado"):
            dar_formato_run(nuevo.findall(qn("a:r"))[0], negritas=True,
                            color=config.COLOR_ENCABEZADO_ESTATUS)
        else:
            aplicar_sangria(nuevo, config.SANGRIA_PROYECTO)
            if renglon.get("sufijo"):
                clonar_run(nuevo, f" {renglon['sufijo']}", negritas=True,
                           color=config.COLOR_SUFIJO_CONCLUIDO)

        cuerpo.append(nuevo)


def escribir_celda(tabla, fila, columna, valor):
    """Escribe un valor en una celda de tabla conservando su formato.
    Parametros:
    tabla: objeto table de python-pptx
    fila: indice de fila
    columna: indice de columna
    valor: texto a escribir, o None para dejar el placeholder"""
    if valor is None:
        return
    celda = tabla.cell(fila, columna)
    parrafos = celda.text_frame.paragraphs
    if not parrafos:
        return
    escribir_valor(parrafos[0]._p, valor)


def agregar_nota_jornadas(slide, texto):
    """Crea el textbox de la nota al pie de Jornadas de Paz.
    Parametros:
    slide: diapositiva de la ficha
    texto: contenido de la nota"""
    caja = slide.shapes.add_textbox(
        Cm(config.NOTA_JP_LEFT), Cm(config.NOTA_JP_TOP),
        Cm(config.NOTA_JP_ANCHO), Cm(config.NOTA_JP_ALTO),
    )
    marco = caja.text_frame
    marco.word_wrap = True
    parrafo = marco.paragraphs[0]
    parrafo.text = texto
    for run in parrafo.runs:
        run.font.name = config.NOTA_JP_FUENTE
        run.font.size = Pt(config.NOTA_JP_TAMANO)
        run.font.color.rgb = RGBColor(0, 0, 0)
    return caja


def escribir_bloque_listas(shapes, claves, renglones_1, renglones_2):
    """Escribe las dos columnas de un listado y elimina las que quedan vacias.
    Parametros:
    shapes: diccionario de shapes por nombre
    claves: diccionario de nombres de shape del bloque
    renglones_1: renglones de la primera columna
    renglones_2: renglones de la segunda columna"""
    shape_1 = shapes.get(claves["lista_1"])
    shape_2 = shapes.get(claves["lista_2"])

    if renglones_1:
        escribir_listado(shape_1, renglones_1)
    elif shape_1 is not None:
        eliminar_shape(shape_1)

    if renglones_2:
        escribir_listado(shape_2, renglones_2)
    elif shape_2 is not None:
        eliminar_shape(shape_2)


def reacomodar_ceci(shapes, claves_visibles):
    """Reparte los grupos visibles de CECI a lo ancho del area, centrando cada
    uno en su porcion. Se usa cuando se elimina la casilla de concluido.
    Parametros:
    shapes: diccionario de shapes por nombre
    claves_visibles: llaves de GRUPOS_CECI que siguen en la diapositiva"""
    grupos = [shapes[config.GRUPOS_CECI[c]] for c in claves_visibles
              if config.GRUPOS_CECI.get(c) in shapes]
    if not grupos:
        return

    izquierda = Cm(config.AREA_CECI["left"])
    ancho = Cm(config.AREA_CECI["ancho"])
    porcion = ancho / len(grupos)

    for posicion, grupo in enumerate(grupos):
        centro = izquierda + porcion * (posicion + 0.5)
        grupo.left = int(centro - grupo.width / 2)


def renderizar(ficha, ruta_plantilla, ruta_salida, repartir_listado):
    """Genera el pptx de una entidad a partir de su FichaData.
    Parametros:
    ficha: objeto FichaData ya formateado
    ruta_plantilla: ruta del pptx de plantilla
    ruta_salida: ruta del archivo a guardar
    repartir_listado: funcion que divide un listado en dos columnas"""
    prs = Presentation(ruta_plantilla)
    slide = prs.slides[0]
    shapes = recolectar_shapes(slide.shapes)

    escribir_texto(shapes.get(config.SHAPE_ENTIDAD), ficha.entidad)

    for llave, nombre_shape in config.SHAPES_INVERSION.items():
        escribir_texto(shapes.get(nombre_shape), ficha.inversion.get(llave))

    for llave, nombre_shape in config.SHAPES_HOSPITALES.items():
        if llave.startswith("lista"):
            continue
        escribir_texto(shapes.get(nombre_shape), ficha.hospitales.get(llave))

    for llave, nombre_shape in config.SHAPES_UMF.items():
        if llave.startswith("lista"):
            continue
        escribir_texto(shapes.get(nombre_shape), ficha.umf.get(llave))

    col_1, col_2 = repartir_listado(ficha.lista_hospitales)
    escribir_bloque_listas(shapes, config.SHAPES_HOSPITALES, col_1, col_2)

    col_1, col_2 = repartir_listado(ficha.lista_umf)
    escribir_bloque_listas(shapes, config.SHAPES_UMF, col_1, col_2)

    tabla_shape = shapes.get(config.SHAPE_TABLA_EQUIPO)
    if tabla_shape is not None and tabla_shape.has_table:
        tabla = tabla_shape.table
        for fila, llave in enumerate(config.ORDEN_TABLA_EQUIPO):
            if fila < len(tabla.rows):
                escribir_celda(tabla, fila, config.COLUMNA_CIFRA_TABLA_EQUIPO,
                               ficha.equipo.get(llave))

    escribir_texto(shapes.get(config.SHAPE_EQ_MDP), ficha.equipo_mdp)
    escribir_texto(shapes.get(config.SHAPE_DRAFT), ficha.draft)

    # CECI. La casilla de concluido solo aparece en las entidades que ya tienen
    # alguno; en las demas se elimina el grupo y se reacomodan las restantes.
    claves_ceci = list(config.GRUPOS_CECI)
    if not ficha.aplica_ceci_concluido:
        grupo = shapes.get(config.GRUPOS_CECI["concluido"])
        if grupo is not None:
            eliminar_shape(grupo)
        claves_ceci.remove("concluido")

    for llave, nombre_shape in config.SHAPES_CECI.items():
        if llave == "concluido" and not ficha.aplica_ceci_concluido:
            continue
        escribir_texto(shapes.get(nombre_shape), ficha.ceci.get(llave))

    reacomodar_ceci(shapes, claves_ceci)

    for llave, nombre_shape in config.SHAPES_VIVE_SALUDABLE.items():
        escribir_texto(shapes.get(nombre_shape), ficha.vive_saludable.get(llave))

    for llave, nombre_shape in config.SHAPES_CASA_X_CASA.items():
        escribir_texto(shapes.get(nombre_shape), ficha.casa_x_casa.get(llave))

    # Jornadas de Paz. El bloque se elimina completo cuando la entidad no
    # aparece en el archivo.
    if ficha.aplica_jornadas_paz:
        for llave, nombre_shape in config.SHAPES_JP.items():
            escribir_texto(shapes.get(nombre_shape), ficha.jornadas_paz.get(llave))
        if ficha.nota_jornadas_paz:
            agregar_nota_jornadas(slide, ficha.nota_jornadas_paz)
    else:
        grupo = shapes.get(config.SHAPE_GRUPO_JP)
        if grupo is not None:
            eliminar_shape(grupo)

    # Mexico Te Abraza, mismo criterio.
    if ficha.aplica_mexico_te_abraza:
        for llave, nombre_shape in config.SHAPES_MTA.items():
            escribir_texto(shapes.get(nombre_shape), ficha.mexico_te_abraza.get(llave))
    else:
        grupo = shapes.get(config.SHAPE_GRUPO_MTA)
        if grupo is not None:
            eliminar_shape(grupo)

    prs.save(ruta_salida)
    return ruta_salida